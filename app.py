from flask import Flask, request, send_file
from pptx import Presentation
from io import BytesIO

app = Flask(__name__)

@app.route("/generate", methods=["POST"])
def generate():
    data = request.json
    slides = data.get("slides", [])

    prs = Presentation()

    for slide_data in slides:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_data.get("title", "")
        content = slide.placeholders[1]
        content.text = "\n".join(slide_data.get("bullets", []))

    output = BytesIO()
    prs.save(output)
    output.seek(0)

    return send_file(
        output,
        as_attachment=True,
        download_name="reporte.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

if __name__ == "__main__":
    app.run()
